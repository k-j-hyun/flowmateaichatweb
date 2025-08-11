import os
from pptx import Presentation
from io import BytesIO, StringIO
from PIL import Image
from utils.image_utils import analyze_image_with_qwen

def pptx_to_markdown_string(pptx_path: str, mode:str = "simple") -> str:
    prs = Presentation(pptx_path)
    output = StringIO()

    image_output_dir = "temp_imgs"
    os.makedirs(image_output_dir, exist_ok=True)

    output.write("# PPT 자동 변환\n\n")

    img_idx = 1

    for i, slide in enumerate(prs.slides):
        output.write(f"## 슬라이드 {i + 1}\n\n")

        for shape in slide.shapes:
            # 텍스트 추출
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    output.write(f"{text}\n\n")

            # 이미지 추출 및 분석
            if shape.shape_type == 13:  # picture
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext if image.ext else "png"
                image_path = os.path.join(image_output_dir, f"slide_{i+1}_img_{img_idx}.{image_ext}")

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                # 이미지 분석
                try:
                    result = analyze_image_with_qwen(image_path, mode=mode)
                    output.write(f"**[이미지 {img_idx} 분석 결과]**\n{result.strip()}\n\n")
                except Exception as e:
                    output.write(f"[이미지 분석 실패: {e}]\n\n")

                img_idx += 1

        output.write("---\n\n")

    return output.getvalue()


if __name__ == "__main__":
    FILE_PATH = "sample_inputs/sample.pptx"
    result = pptx_to_markdown_string(FILE_PATH)
    print(result)
