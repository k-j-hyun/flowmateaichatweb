"""
PowerPoint 프레젠테이션 생성 모듈
마크다운 텍스트를 전문적인 PPT로 변환
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from langchain_ollama import ChatOllama
import re
import os
import logging

# 로깅 설정
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTXGenerator:
    """PowerPoint 생성기 클래스"""
    
    def __init__(self):
        """초기화"""
        self.prs = None
        self.theme_colors = {
            'primary': RGBColor(99, 102, 241),    # 보라색
            'secondary': RGBColor(139, 92, 246),  # 연한 보라색
            'accent': RGBColor(124, 58, 237),     # 진한 보라색
            'text': RGBColor(74, 85, 104),        # 회색
            'white': RGBColor(255, 255, 255),     # 흰색
            'background': RGBColor(248, 250, 252) # 연한 회색
        }
    
    def create_presentation(self):
        """새 프레젠테이션 생성"""
        try:
            self.prs = Presentation()
            # 16:9 비율 설정
            self.prs.slide_width = Inches(13.33)
            self.prs.slide_height = Inches(7.5)
            logger.info("프레젠테이션 객체 생성 완료")
            return True
        except Exception as e:
            logger.error(f"프레젠테이션 생성 실패: {e}")
            return False
    
    def add_title_slide(self, title, subtitle=""):
        """제목 슬라이드 추가"""
        try:
            title_layout = self.prs.slide_layouts[0]
            slide = self.prs.slides.add_slide(title_layout)
            
            # 제목 설정
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.alignment = PP_ALIGN.CENTER
                title_paragraph.font.size = Pt(44)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = self.theme_colors['primary']
            
            # 부제목 설정
            if subtitle and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle
                subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
                subtitle_paragraph.alignment = PP_ALIGN.CENTER
                subtitle_paragraph.font.size = Pt(22)
                subtitle_paragraph.font.color.rgb = self.theme_colors['text']
            
            logger.info(f"제목 슬라이드 추가: {title}")
            return slide
        except Exception as e:
            logger.error(f"제목 슬라이드 추가 실패: {e}")
            return None
    
    def add_content_slide(self, title, points, slide_type="bullet"):
        """내용 슬라이드 추가"""
        try:
            if slide_type == "two_column" and len(points) > 5:
                layout = self.prs.slide_layouts[3]  # 두 컬럼 레이아웃
            else:
                layout = self.prs.slide_layouts[1]  # 제목 + 내용 레이아웃
            
            slide = self.prs.slides.add_slide(layout)
            
            # 제목 설정
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title
                title_shape.left = Inches(0.5)
                title_shape.top = Inches(0.3)
                title_shape.width = Inches(12.3)
                title_shape.height = Inches(1.2)
                
                title_paragraph = title_shape.text_frame.paragraphs[0]
                title_paragraph.font.size = Pt(36)
                title_paragraph.font.bold = True
                title_paragraph.font.color.rgb = self.theme_colors['primary']
                
                # 배경 색상 설정
                fill = title_shape.fill
                fill.solid()
                fill.fore_color.rgb = self.theme_colors['background']
            
            # 내용 추가
            if slide_type == "two_column" and len(points) > 5:
                self._add_two_column_content(slide, points)
            else:
                self._add_single_column_content(slide, points)
            
            logger.info(f"내용 슬라이드 추가: {title}")
            return slide
        except Exception as e:
            logger.error(f"내용 슬라이드 추가 실패: {e}")
            return None
    
    def _add_single_column_content(self, slide, points):
        """단일 컬럼 내용 추가"""
        try:
            if len(slide.placeholders) > 1:
                content_frame = slide.placeholders[1].text_frame
                content_frame.clear()
                
                icons = ["-"] * 5
                
                for i, point in enumerate(points):
                    p = content_frame.add_paragraph()
                    p.text = f"{icons[i % len(icons)]} {point}"
                    p.level = 0
                    p.font.size = Pt(22)
                    p.space_after = Pt(15)
        except Exception as e:
            logger.error(f"단일 컬럼 내용 추가 실패: {e}")
    
    def _add_two_column_content(self, slide, points):
        """두 컬럼 내용 추가"""
        try:
            if len(slide.placeholders) > 2:
                left_content = slide.placeholders[1].text_frame
                right_content = slide.placeholders[2].text_frame
                
                mid_point = len(points) // 2
                
                # 왼쪽 컬럼
                left_content.clear()
                for point in points[:mid_point]:
                    p = left_content.add_paragraph()
                    p.text = f"• {point}"
                    p.level = 0
                    p.font.size = Pt(20)
                    p.space_after = Pt(12)
                
                # 오른쪽 컬럼
                right_content.clear()
                for point in points[mid_point:]:
                    p = right_content.add_paragraph()
                    p.text = f"• {point}"
                    p.level = 0
                    p.font.size = Pt(20)
                    p.space_after = Pt(12)
        except Exception as e:
            logger.error(f"두 컬럼 내용 추가 실패: {e}")
    
    def add_table_slide(self, title, table_data):
        """테이블 슬라이드 추가"""
        try:
            layout = self.prs.slide_layouts[5]  # 빈 레이아웃
            slide = self.prs.slides.add_slide(layout)
            
            # 제목 추가
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = title
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(24)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = self.theme_colors['primary']
            
            # 테이블 생성
            if table_data and len(table_data) > 0:
                rows = len(table_data)
                cols = len(table_data[0]) if table_data[0] else 1
                
                table = slide.shapes.add_table(
                    rows, cols, Inches(1), Inches(2), Inches(11), Inches(4.5)
                ).table
                
                # 테이블 데이터 입력
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_data in enumerate(row_data):
                        if col_idx < cols:
                            cell = table.cell(row_idx, col_idx)
                            cell.text = str(cell_data)
                            
                            # 헤더 행 스타일링
                            if row_idx == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = self.theme_colors['primary']
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.color.rgb = self.theme_colors['white']
                                        run.font.bold = True
                                        run.font.size = Pt(14)
                            else:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
            
            logger.info(f"테이블 슬라이드 추가: {title}")
            return slide
        except Exception as e:
            logger.error(f"테이블 슬라이드 추가 실패: {e}")
            return None
    
    def save_presentation(self, output_path):
        """프레젠테이션 저장"""
        try:
            # 출력 디렉토리 생성
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            self.prs.save(output_path)
            logger.info(f"프레젠테이션 저장 완료: {output_path}")
            return True
        except Exception as e:
            logger.error(f"프레젠테이션 저장 실패: {e}")
            return False

def generate_slide_structure(text_content):
    """LLM을 사용하여 슬라이드 구조 생성"""
    try:
        llm = ChatOllama(model="anpigon/qwen2.5-7b-instruct-kowiki:latest")
        
        prompt = f"""
당신은 한국어 발표 슬라이드 전문가입니다.
다음 텍스트를 바탕으로 PPT 슬라이드 구조를 만들어주세요.
반드시 한국어로 응답하세요.

입력 텍스트:
{text_content}

출력 형식:
[슬라이드 1]
제목: 프로젝트 소개
핵심 포인트:
- 첫 번째 포인트
- 두 번째 포인트

[슬라이드 2]
제목: 주요 내용
핵심 포인트:
- 첫 번째 포인트
- 두 번째 포인트

각 슬라이드는 3-5개의 핵심 포인트를 가져야 합니다.
제목은 명확하고 간결해야 합니다.
"""
        
        response = llm.invoke([prompt])
        logger.info("슬라이드 구조 생성 완료")
        return response.content
    except Exception as e:
        logger.error(f"슬라이드 구조 생성 실패: {e}")
        return None

def parse_slide_structure(structured_text):
    """구조화된 텍스트를 파싱하여 슬라이드 데이터 추출 - qwen 모델 최적화"""
    try:
        slides = []
        logger.info(f"파싱할 텍스트: {structured_text[:500]}...")  # 디버깅용 로그
        
        # qwen 모델의 출력 패턴에 맞춘 슬라이드 구분자
        slide_patterns = [
            r"\[슬라이드 \d+\]",
            r"슬라이드 \d+:",
            r"### 슬라이드 \d+",
            r"## 슬라이드 \d+",
            r"\d+\.\s*슬라이드"
        ]
        
        # 가장 적합한 패턴 찾기
        best_pattern = None
        for pattern in slide_patterns:
            matches = re.findall(pattern, structured_text)
            if matches:
                best_pattern = pattern
                logger.info(f"사용 패턴: {pattern}, 찾은 매치: {len(matches)}개")
                break
        
        if not best_pattern:
            # 패턴이 없으면 폴백 처리
            logger.warning("슬라이드 구분자를 찾을 수 없음, 폴백 모드 사용")
            return create_fallback_slides(structured_text)
        
        # 슬라이드 블록 분할
        slide_blocks = re.split(best_pattern, structured_text)
        slide_blocks = [block.strip() for block in slide_blocks if block.strip()]
        
        for i, block in enumerate(slide_blocks):
            # qwen 모델 출력에 최적화된 제목 추출
            title_patterns = [
                r"제목:\s*(.+)",
                r"\*\*제목\*\*:\s*(.+)",
                r"제목\s*:\s*(.+)",
                r"^(.+?)(?:\n|$)"  # 첫 번째 줄을 제목으로
            ]
            
            title = f"슬라이드 {i+1}"  # 기본 제목
            for pattern in title_patterns:
                title_match = re.search(pattern, block, re.MULTILINE | re.DOTALL)
                if title_match:
                    candidate_title = title_match.group(1).strip()
                    
                    # "제목:" 접두사 제거
                    if candidate_title.startswith('제목:'):
                        candidate_title = candidate_title[3:].strip()
                    elif candidate_title.startswith('제목 :'):
                        candidate_title = candidate_title[4:].strip()
                    
                    # 제목에 '핵심 포인트' 또는 '핵심'이 포함되지 않고, 적절한 길이인 경우에만 사용
                    if (len(candidate_title) < 100 and 
                        not candidate_title.startswith('핵심') and 
                        '핵심 포인트' not in candidate_title and
                        '핵심포인트' not in candidate_title and
                        candidate_title):  # 빈 문자열이 아닌 경우만
                        title = candidate_title
                        break
            
            # qwen 모델에 최적화된 핵심 포인트 추출
            bullet_patterns = [
                r"핵심 포인트:\s*\n((?:[-*•]\s+.+\n?)+)",  # "핵심 포인트:" 이후의 불릿들
                r"[-*•]\s+(.+)",  # 일반 불릿 포인트
                r"\d+\.\s+(.+)",  # 숫자 리스트
                r"▶\s+(.+)",     # 화살표
                r"✓\s+(.+)"      # 체크마크
            ]
            
            points = []
            
            # 먼저 "핵심 포인트:" 섹션을 찾아 처리
            core_points_match = re.search(r"핵심 포인트:\s*\n((?:[-*•]\s+.+(?:\n|$))+)", block, re.MULTILINE)
            if core_points_match:
                core_section = core_points_match.group(1)
                points = re.findall(r"[-*•]\s+(.+)", core_section)
            else:
                # 일반적인 불릿 포인트 패턴 사용
                for pattern in bullet_patterns[1:]:  # 첫 번째 패턴 제외
                    found_points = re.findall(pattern, block, re.MULTILINE)
                    if found_points:
                        points.extend(found_points)
                        break  # 첫 번째로 찾은 패턴만 사용
            
            # 포인트 정리 및 필터링
            cleaned_points = []
            for point in points:
                cleaned_point = point.strip()
                # 너무 짧거나 의미없는 포인트 제외
                if len(cleaned_point) > 3 and not cleaned_point.startswith('제목'):
                    cleaned_points.append(cleaned_point)
            
            # 포인트가 여전히 없으면 문장 단위로 분할
            if not cleaned_points:
                sentences = [s.strip() for s in block.split('\n') 
                           if s.strip() and len(s.strip()) > 3 and '제목:' not in s and '핵심 포인트:' not in s]
                cleaned_points = sentences[:4]  # 최대 4개 포인트
            
            # 테이블 데이터 검색 (qwen 모델의 테이블 형식)
            table_pattern = r"\|(.+?)\|"
            table_matches = re.findall(table_pattern, block, re.MULTILINE)
            
            slide_data = {
                'title': title,
                'points': cleaned_points,
                'table_data': None,
                'slide_type': 'content'
            }
            
            # 테이블 데이터가 있으면 파싱
            if table_matches and len(table_matches) > 1:
                table_data = []
                for match in table_matches:
                    row = [cell.strip() for cell in match.split('|') if cell.strip()]
                    if row:
                        table_data.append(row)
                slide_data['table_data'] = table_data
                slide_data['slide_type'] = 'table'
            
            # 포인트 수에 따라 레이아웃 결정
            elif len(cleaned_points) > 5:
                slide_data['slide_type'] = 'two_column'
            
            # 유효한 슬라이드만 추가
            if slide_data['title'] and (slide_data['points'] or slide_data['table_data']):
                slides.append(slide_data)
        
        logger.info(f"슬라이드 파싱 완료: {len(slides)}개 슬라이드")
        for i, slide in enumerate(slides):
            logger.info(f"슬라이드 {i+1}: {slide['title']} ({len(slide['points'])}개 포인트)")
        
        return slides
    except Exception as e:
        logger.error(f"슬라이드 파싱 실패: {e}")
        return []

def create_presentation_from_text(text_content, output_path="output/presentation.pptx"):
    """텍스트 내용으로부터 프레젠테이션 생성 - qwen 모델 최적화"""
    try:
        logger.info("프레젠테이션 생성 시작")
        
        # PPT 생성기 초기화
        generator = PPTXGenerator()
        if not generator.create_presentation():
            return None
        
        # LLM으로 슬라이드 구조 생성
        # structured_text = generate_slide_structure(text_content)
        # 이미 views.py에서 실행함
        structured_text = text_content
        if not structured_text:
            logger.error("슬라이드 구조 생성 실패")
            return None
        
        # 슬라이드 데이터 파싱
        slides_data = parse_slide_structure(structured_text)
        
        # 파싱 실패 시 폴백: 간단한 슬라이드 생성
        if not slides_data:
            logger.warning("슬라이드 파싱 실패, 폴백 모드로 간단한 슬라이드 생성")
            slides_data = create_fallback_slides(structured_text)
        
        if not slides_data:
            logger.error("폴백 슬라이드 생성도 실패")
            return None
        
        # 제목 슬라이드 추가
        if slides_data:
            main_title = slides_data[0]['title'] if slides_data[0]['title'] else "발표 자료"
            generator.add_title_slide(main_title, "FlowMate AI 발표자료")
        
        # 내용 슬라이드 추가
        for slide_data in slides_data:
            try:
                if slide_data['slide_type'] == 'table' and slide_data['table_data']:
                    generator.add_table_slide(slide_data['title'], slide_data['table_data'])
                else:
                    slide_type = slide_data['slide_type'] if slide_data['slide_type'] != 'content' else 'bullet'
                    generator.add_content_slide(slide_data['title'], slide_data['points'], slide_type)
            except Exception as slide_error:
                logger.warning(f"슬라이드 추가 실패: {slide_error}, 건너뜀")
                continue
        
        # 감사 슬라이드 추가
        generator.add_title_slide("감사합니다", "질문이 있으시면 언제든 말씀해 주세요")
        
        # 프레젠테이션 저장
        if generator.save_presentation(output_path):
            logger.info(f"프레젠테이션 생성 완료: {output_path}")
            return output_path
        else:
            return None
            
    except Exception as e:
        logger.error(f"프레젠테이션 생성 중 오류: {e}")
        return None

def create_fallback_slides(text_content):
    """파싱 실패 시 폴백용 간단한 슬라이드 생성"""
    try:
        logger.info("폴백 슬라이드 생성 시작")
        
        # 텍스트를 문장 단위로 분할
        sentences = [s.strip() for s in text_content.split('\n') if s.strip()]
        
        if not sentences:
            sentences = [s.strip() for s in text_content.split('.') if s.strip()]
        
        slides = []
        
        # 첫 번째 슬라이드: 개요
        overview_points = sentences[:5] if len(sentences) > 5 else sentences
        slides.append({
            'title': '개요',
            'points': overview_points,
            'table_data': None,
            'slide_type': 'content'
        })
        
        # 나머지 내용을 여러 슬라이드로 분할
        remaining_sentences = sentences[5:] if len(sentences) > 5 else []
        
        slide_count = 2
        chunk_size = 4  # 슬라이드당 4개 포인트
        
        for i in range(0, len(remaining_sentences), chunk_size):
            chunk = remaining_sentences[i:i+chunk_size]
            if chunk:
                slides.append({
                    'title': f'주요 내용 {slide_count - 1}',
                    'points': chunk,
                    'table_data': None,
                    'slide_type': 'content'
                })
                slide_count += 1
        
        # 최소 1개 슬라이드는 보장
        if not slides:
            slides.append({
                'title': '내용',
                'points': ['내용을 처리하는 중 문제가 발생했습니다.', '원본 텍스트를 확인해 주세요.'],
                'table_data': None,
                'slide_type': 'content'
            })
        
        logger.info(f"폴백 슬라이드 생성 완료: {len(slides)}개")
        return slides
        
    except Exception as e:
        logger.error(f"폴백 슬라이드 생성 실패: {e}")
        return []

# 하위 호환성을 위한 함수
def save_structured_text_to_pptx(whole_text, output_path="output/presentation.pptx"):
    """기존 함수명 유지 (하위 호환성)"""
    return create_presentation_from_text(whole_text, output_path)

def make_text_to_slide_text(whole_text):
    """기존 함수명 유지 (하위 호환성)"""
    return generate_slide_structure(whole_text)

if __name__ == "__main__":
    # 테스트용 샘플 텍스트
    sample_text = """
FlowMate AI는 로컬 환경에서 구동되는 AI 기반 업무 솔루션입니다.
주요 기능으로는 AI 채팅 어시스턴트, 문서 분석 및 요약, 발표 영상 분석, HR 예측 모델이 있습니다.
기술 스택은 Django, Ollama qwen2.5:7b, BGE-M3 임베딩을 사용합니다.
보안 특징으로는 로컬 처리, 외부 전송 없음, On-premise 배포가 가능합니다.
향후 계획으로는 모델 성능 개선, 다국어 지원, API 확장, 모바일 앱 개발이 있습니다.
"""
    
    # 프레젠테이션 생성 테스트
    output_path = create_presentation_from_text(sample_text, "output/test_presentation.pptx")
    
    if output_path:
        print(f"✅ 프레젠테이션 생성 성공: {output_path}")
        print("주요 개선사항:")
        print("- 한국어 전용 처리")
        print("- 예외 처리 강화")
        print("- 로깅 시스템 추가")
        print("- 클래스 기반 구조")
        print("- 깔끔한 코드 정리")
    else:
        print("❌ 프레젠테이션 생성 실패")